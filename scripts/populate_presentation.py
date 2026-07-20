import pptx
import os

replacements = {
    "[Student Name]": "Srijan Basnet",
    "[Exam Roll Number]": "370121",
    "[Registration Number]": "025-3-2-04314-2022",
    "[Date of Defense]": "July 18, 2026",
    "[Program]": "Bachelor of Computer Applications (BCA)",
    "[X]th": "8th",
    "[College Name]": "Kantipur City College",
    "[Programming languages, frameworks, tools]": "• Language: Python (v3.11)\n• NLP & AI/ML: FAISS, Sentence-Transformers, LangChain, PyMuPDF\n• Web & Serving: FastAPI, Uvicorn, Lemonade Server, Ollama\n• Models: Qwen3-1.7B-GGUF, Whisper, Kokoro TTS",
    "[Web development / Data / AI / Networking, etc.]": "Artificial Intelligence, Machine Learning, Natural Language Processing, Conversational Systems",
    "[Short statement on career goals]": "To become a professional AI/ML and NLP developer, building offline, secure, and highly efficient conversational applications using state-of-the-art local LLM technologies.",
    "[Organization Name]": "Research Management Cell (RMC) Department, Kantipur City College",
    "[industry]": "academic research and software development",
    "[year]": "1998",
    "[location]": "Putalisadak, Kathmandu",
    "[core focus]": "AI/ML interfaces and college management solutions",
    "[Mission statement placeholder]": "To foster academic research, facilitate technological innovation, and execute software development projects that serve the academic and administrative requirements of the institution.",
    "[Vision statement placeholder]": "To build an ecosystem that integrates academic teaching with research and development, bridging the gap between theoretical computing and real-world technology.",
    "[Product / Service 1]": "R&D facilitation and academic research publishing",
    "[Product / Service 2]": "Institutional web portals and digital services",
    "[Product / Service 3]": "AI prototypes (e.g. ANA - AI Receptionist Assistant)",
    "[Start]": "May 13, 2026",
    "[End]": "July 8, 2026",
    "[Department Name]": "Research Management Cell (RMC) Department",
    "[Team Name]": "AI & Conversational NLP Team",
    "[Job title]": "AI/ML Intern (Speech & NLP)",
    "[Key responsibilities]": "• Design and implement the Retrieval-Augmented Generation (RAG) pipeline.\n• Build multi-file vector indexes for college documents.\n• Implement domain-specific keyword routing and lexicon normalization.\n• Integrate the local LLM (Qwen3) served via Lemonade server.",
    "[Languages, frameworks,\nsoftware, platforms]": "• Python 3.11, FastAPI\n• FAISS, Sentence-Transformers\n• LangChain, PyMuPDF\n• Lemonade, Ollama",
    "[Languages, frameworks,": "• Python 3.11, FastAPI\n• FAISS, Sentence-Transformers",
    "software, platforms]": "• LangChain, PyMuPDF\n• Lemonade, Ollama (Qwen3)",
    "A structured 6-week engagement": "A structured 8-week engagement combining people, process, and technology",
    "[Photo Placeholder]": "[Photo: Srijan Basnet]",
    "(~12 weeks)": "(8 weeks)",
    "Weekly Work Log (Week 1–3)": "Weekly Work Log (Week 1–3)",
    "Weekly Work Log (Week 3–4)": "Weekly Work Log (Week 4–6)",
    
    # Slide 4 Table (Weeks 1-3)
    "[Onboarding, environment setup]": "Setup Python 3.11 development env, venv, Git. Scaffolded project monorepo structure.",
    "[First assigned mini-task]": "Prepared college text data files. Developed document chunking logic using LangChain text splitters.",
    "[Assigned to specific module/task]": "Built FAISS vector indexes for semantic similarity search on college files.",
    "[Introduced to codebase / tools]": "Familiarized with project monorepo layout and overall architecture of the ANA AI receptionist.",
    "[Understood workflow / SDLC]": "Analyzed overlap and chunk size configurations for retrieval optimization.",
    "[Contributed to feature development]": "Mastered FAISS index creation and querying for multi-file search.",
    "[Adapting to new tools/processes]": "Resolved virtual environment package version conflicts.",
    "[Understanding existing code]": "Handled irregular formatting in raw college source documents.",
    "[Understanding requirements]": "Addressed domain-specific queries and overlap relevance.",
    "[Familiarization with tech stack]": "Environment configuration, Git, package dependency management.",
    "[Debugging, version control]": "Text extraction from PDF files (PyMuPDF) and structured data preprocessing.",
    "[Applied core technical skill]": "Vector database indexing, similarity search configurations.",

    # Slide 5 Table (Weeks 4-6)
    "[Testing / bug fixing]": "Integrated local Qwen3-1.7B-GGUF via Lemonade. Set up Whisper STT and Kokoro TTS voice modules.",
    "[Core project development]": "Implemented Wav2Lip talking-head avatar wrapper. Connected TTS audio with lip-sync generation.",
    "[Finalization & documentation]": "Implemented time-aware schedule handler. Built RAG fallback layer and completed project documentation.",
    "[Learned code review process]": "Understood model serving configurations and local API endpoints.",
    "[Integration with systems]": "Learned model pipeline integration (Whisper + Qwen3 + Kokoro + Wav2Lip).",
    "[Presented work to team]": "Presented the working voice and talking-head RAG interface to RMC supervisors.",
    "[Coordinating with team]": "Optimized latency of Wav2Lip video inference on local CPU/GPU hardware.",
    "[Handling edge cases]": "Resolved Whisper Nepali accent transcription mismatches and local model memory limits.",
    "[Meeting deadlines]": "Managed parallel computation pipelines for video rendering.",
    "[Collaboration & communication]": "FastAPI backend endpoint development, prompt engineering.",
    "[Advanced technical skill]": "Audio-video synchronization, lip-sync rendering, video caching.",
    "[Presentation & documentation]": "Nepal timezone datetime logic, RAG confidence thresholds, and final report writing.",

    # Slide 6
    "[Project Title Placeholder]": "ANA (AI Receptionist Assistant) - Local RAG Pipeline",
    "[Describe the problem the project addressed]": "• Inconsistent formatting in unstructured college data leads to irrelevant retrieval.\n• General-purpose models struggle with local acronyms (BCA-IT) and teacher schedules.\n• Local LLM generation introduces latency, causing conversation to feel sluggish.\n• AI hallucinations on queries where source documents have no direct match.",
    "[Summarize the solution/approach]": "• Implemented multi-file indexing pipeline using Sentence-Transformers and FAISS.\n• Built domain-specific keyword routing and custom lexicon normalization.\n• Integrated offline Qwen3-1.7B-GGUF served via local Lemonade server.\n• Developed a confidence-based fallback mechanism to guarantee relevant answers.",
    "[Languages, frameworks, DB, tools]": "• Python 3.11, FastAPI\n• FAISS (vector database)\n• Sentence-Transformers (all-MiniLM-L6-v2)\n• LangChain, PyMuPDF, Lemonade / Ollama",
    "[Your specific role and impact]": "• Designed & built the end-to-end local RAG indexing and search module.\n• Developed the keyword routing and custom normalization layer to handle KCC terms.\n• Created FastAPI backend routes and integrated speech-to-text / text-to-speech APIs.",
    "[Insert project screenshot here]": "Local RAG Pipeline & Multi-File Indexing Architecture",

    # Slide 7
    "[Language / framework]": "Python 3.11 & FastAPI endpoints",
    "[Tool / platform]": "Ollama & Lemonade local LLM serving",
    "[Database / API]": "FAISS vector search database",
    "[Testing / debugging]": "Retrieval accuracy and confidence thresholds",
    "[Technical challenge]": "Local LLM inference latency on consumer CPU, optimized via model quantization.",
    "[Time management]": "Synchronizing parallel speech-to-text, text-to-speech, and video rendering pipelines.",
    "[Communication in team]": "Integrating NLP services with the frontend React interface and speech modules.",
    "[Adapting to new tools]": "Learning local LLM parameters, GGUF quantization, and vector indexing.",
    "[Professional work culture]": "Adapting to Research Management Cell (RMC) documentation and development standards.",
    "[Real-world application of theory]": "Translating NLP and vector space theories into offline, low-latency applications.",
    "[Teamwork & collaboration]": "Cooperating with speech and avatar development interns under internal supervision.",
    "[Problem-solving mindset]": "Debugging local GPU memory limits and building a high-efficiency video cache.",

    # Slide 8
    "[Brief wrap-up of overall experience]": "• Developed a fully offline, high-performance RAG pipeline for KCC's AI Receptionist.\n• Demonstrated the feasibility of hosting lightweight open-source models on consumer hardware.\n• Achieved low-latency conversational responses without compromising user privacy or incurring API costs.",
    "[What could be done better / next steps]": "• Integrate vector reranking models (e.g. cross-encoders) to improve top-k similarity precision.\n• Implement query expansion (HYDE) to handle highly abbreviated student queries.\n• Set up an active feedback loop to capture user queries and expand the text document index.",
    "[How this shaped your career direction]": "• Solidified interest in Conversational AI and Retrieval architectures.\n• Built deep expertise in vector databases, speech processing, and local LLM pipelines.",
}

def replace_in_text_frame(tf):
    # Try whole-text replacement for multiline placeholders first
    full_text = "\n".join([p.text for p in tf.paragraphs])
    
    # Clean up Windows-1252 non-breaking dash character \x96 (en-dash) to standard en-dash \u2013
    full_text_clean = full_text.replace("\x96", "–")
    
    replaced = False
    for k, v in replacements.items():
        if k in full_text_clean:
            full_text_clean = full_text_clean.replace(k, v)
            replaced = True
            
    if replaced:
        tf.clear()
        lines = full_text_clean.split("\n")
        for idx, line in enumerate(lines):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
        return

    # Fallback to run-by-run replacement for single-line placeholders
    for p in tf.paragraphs:
        for run in p.runs:
            run_text_clean = run.text.replace("\x96", "–")
            replaced_run = False
            for k, v in replacements.items():
                if k in run_text_clean:
                    run_text_clean = run_text_clean.replace(k, v)
                    replaced_run = True
            if replaced_run:
                run.text = run_text_clean

def replace_in_table(table):
    for row in table.rows:
        for cell in row.cells:
            replace_in_text_frame(cell.text_frame)

def restore_and_update_presentation(filepath):
    # If the file has already been partially updated, some keys might have been lost
    # We can check if the file exists and run update.
    # In order to make it robust, we will just run the replacement.
    if not os.path.exists(filepath):
        print(f"File not found: {filepath}")
        return False
        
    print(f"Updating {filepath}...")
    prs = pptx.Presentation(filepath)
    
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # Check for standard text frames
            if hasattr(shape, "text_frame") and shape.text_frame:
                replace_in_text_frame(shape.text_frame)
                
            # Check for tables
            if shape.has_table:
                replace_in_table(shape.table)
                
    prs.save(filepath)
    print(f"Successfully saved {filepath}\n")
    return True

if __name__ == "__main__":
    paths = [
        "C:/Users/Acer/Desktop/ml_chatbot/Internship_Defense_Presentation.pptx",
        "C:/Users/Acer/Downloads/Internship_Defense_Presentation.pptx"
    ]
    for p in paths:
        restore_and_update_presentation(p)
